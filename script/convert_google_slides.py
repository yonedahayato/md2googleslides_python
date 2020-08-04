from pptx import Presentation

def convert_slide(text):
    """ convert_google_slides func

    convert to slide from text

    Args:
        text(str): text want to convert

    """
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    prs.save('test.pptx')
